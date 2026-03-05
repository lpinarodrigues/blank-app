import streamlit as st
from groq import Groq
import json
import io
import re
import PyPDF2
import docx
import pptx
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.pagesizes import A4
from reportlab.lib.enums import TA_JUSTIFY
from database import salvar_item_estudo

def get_ai_client():
    return Groq(api_key=st.secrets["GROQ_API_KEY"])

def consultar_core_ia_perfeicao(prompt, contexto="", img_b64=None):
    client = get_ai_client()
    
    # Se houver imagem, usamos o modelo Vision. Se não, o modelo de texto avançado.
    modelo = "llama-3.2-11b-vision-preview" if img_b64 else "llama-3.3-70b-versatile"
    
    sys_msg = (
        "Você é o ORÁCULO NEXUS. Sua resposta deve ser de nível PhD em Medicina. "
        "Se houver imagem, analise-a clinicamente (ECG, Lab, Imagem). "
        "Ignore ruídos de OCR. Estruture com markdown: # Título, ## Subtítulo, **Negrito**."
    )
    
    conteudo_msg = []
    if img_b64:
        conteudo_msg.append({"type": "text", "text": f"Contexto: {contexto}\n\nPergunta: {prompt}"})
        conteudo_msg.append({"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{img_b64}"}})
    else:
        conteudo_msg.append({"type": "text", "text": f"Contexto: {contexto}\n\nPergunta: {prompt}"})

    completion = client.chat.completions.create(
        model=modelo,
        messages=[{"role": "system", "content": sys_msg}, {"role": "user", "content": conteudo_msg}],
        temperature=0.2
    )
    return completion.choices[0].message.content, "Clínica Médica", "Geral"

def limpar_json(texto):
    texto = texto.strip()
    m = "```json"; ms = "```"
    if texto.startswith(m): texto = texto[7:]
    elif texto.startswith(ms): texto = texto[3:]
    if texto.endswith(ms): texto = texto[:-3]
    return texto.strip()

def gerar_apenas_flashcards(texto_base, area, subtema, email, qtd=5):
    client = get_ai_client()
    prompt = f"Crie {qtd} Flashcards médicos de elite em JSON: {{'flashcards': [{{'pergunta':'', 'resposta':''}}]}}. Texto: {texto_base}"
    try:
        res = client.chat.completions.create(model="llama-3.3-70b-versatile", messages=[{"role": "user", "content": prompt}], response_format={"type": "json_object"})
        dados = json.loads(limpar_json(res.choices[0].message.content))
        for c in dados.get('flashcards', []):
            salvar_item_estudo({"pergunta": c['pergunta'], "resposta": c['resposta'], "grande_area": area, "subtema": subtema, "categoria": "Flashcard", "criado_por_email": email})
        return len(dados.get('flashcards', []))
    except: return 0

def gerar_apenas_questoes(texto_base, area, subtema, email, qtd=3):
    client = get_ai_client()
    prompt = f"Crie {qtd} questões de múltipla escolha em JSON: {{'questoes': [{{'pergunta':'', 'a':'', 'b':'', 'c':'', 'd':'', 'gabarito':'A', 'explicacao':''}}]}}. Texto: {texto_base}"
    try:
        res = client.chat.completions.create(model="llama-3.3-70b-versatile", messages=[{"role": "user", "content": prompt}], response_format={"type": "json_object"})
        dados = json.loads(limpar_json(res.choices[0].message.content))
        from database import get_supabase
        for q in dados.get('questoes', []):
            get_supabase().table("questionarios").insert({"pergunta": q['pergunta'], "opcao_a": q['a'], "opcao_b": q['b'], "opcao_c": q['c'], "opcao_d": q['d'], "gabarito": q['gabarito'], "explica_correta": q['explicacao'], "criado_por_email": email, "categoria": area}).execute()
        return len(dados.get('questoes', []))
    except: return 0

def ler_documento_universal(arquivo):
    if not arquivo: return ""
    n = arquivo.name.lower()
    try:
        if n.endswith('.pdf'):
            p = PyPDF2.PdfReader(arquivo)
            return " ".join([page.extract_text() for page in p.pages if page.extract_text()])
        elif n.endswith('.docx'):
            d = docx.Document(arquivo)
            return " ".join([p.text for p in d.paragraphs])
        elif n.endswith('.pptx'):
            pt = pptx.Presentation(arquivo)
            return " ".join([shape.text for s in pt.slides for shape in s.shapes if hasattr(shape, "text")])
        return arquivo.getvalue().decode("utf-8", errors="ignore")
    except: return ""

def gerar_pdf_premium(texto, titulo="PROTOCOLO NEXUS CORE"):
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4, rightMargin=50, leftMargin=50, topMargin=50, bottomMargin=50)
    styles = getSampleStyleSheet()
    
    # Estilos Customizados
    estilo_corpo = ParagraphStyle('Corpo', parent=styles['Normal'], fontSize=11, leading=16, alignment=TA_JUSTIFY, spaceAfter=12)
    estilo_titulo = ParagraphStyle('Titulo', parent=styles['Heading1'], fontSize=18, textColor='#1e293b', spaceAfter=20)
    
    story = []
    story.append(Paragraph(f"<b>{titulo.upper()}</b>", estilo_titulo))
    story.append(Spacer(1, 12))
    
    # Limpeza básica de Markdown para ReportLab
    t = texto.replace("**", "<b>").replace("**", "</b>").replace("\n", "<br/>")
    t = re.sub(r'# (.*?)<br/>', r'<h1>\1</h1>', t)
    
    for p in t.split('<br/>'):
        if p.strip(): story.append(Paragraph(p, estilo_corpo))
    
    doc.build(story)
    buffer.seek(0)
    return buffer.getvalue()
